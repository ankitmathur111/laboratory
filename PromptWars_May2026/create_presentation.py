import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    workspace_dir = r"D:\github\laboratory\PromptWars_May2026"
    output_path = os.path.join(workspace_dir, "VoyageFlow_Presentation.pptx")
    
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Premium Color Palette
    NAVY = RGBColor(11, 27, 61)       # Primary #0B1B3D
    GOLD = RGBColor(255, 153, 51)     # Accent #FF9933
    WHITE = RGBColor(255, 255, 255)   # Background
    LIGHT_GRAY = RGBColor(248, 250, 252) # Card background #F8FAFC
    DARK_GRAY = RGBColor(51, 65, 85)   # Text body color #334155
    ACCENT_BLUE = RGBColor(30, 58, 138) # Deep blue for titles #1E3A8A
    
    # Helper to add full background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to create a premium card
    def add_card(slide, left, top, width, height, bg_color=LIGHT_GRAY, border_color=None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background() # No border
        return shape

    # Helper to add text boxes easily
    def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = align
        return txBox

    blank_slide_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Dark Navy Theme)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, NAVY)
    
    # Cover image insertion (inset on the right or used elegantly)
    cover_img_path = os.path.join(workspace_dir, "voyageflow_cover.png")
    if os.path.exists(cover_img_path):
        slide1.shapes.add_picture(cover_img_path, Inches(7.5), Inches(1.0), Inches(5.0), Inches(5.5))
        
    # Title Text Box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "VOYAGEFLOW"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = 'Arial'
    
    p2 = tf1.add_paragraph()
    p2.text = "Dynamic AI Travel Planner"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    p2.font.name = 'Arial'
    p2.space_before = Pt(10)
    
    p3 = tf1.add_paragraph()
    p3.text = "A Premium Spatiotemporal Constraint Solver & Real-Time Event Simulator"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(200, 210, 230)
    p3.font.name = 'Arial'
    p3.space_before = Pt(20)
    
    p4 = tf1.add_paragraph()
    p4.text = "Hosted fully Serverless in WebAssembly on Vercel Edge"
    p4.font.size = Pt(14)
    p4.font.italic = True
    p4.font.color.rgb = RGBColor(160, 180, 210)
    p4.font.name = 'Arial'
    p4.space_before = Pt(40)

    # =========================================================================
    # SLIDE 2: THE VISION (Light Theme + Image)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, WHITE)
    
    # Title
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8), 
                "Dynamic AI Travel Re-invented", font_size=28, bold=True, color=ACCENT_BLUE)
    
    # Subtitle
    add_textbox(slide2, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4), 
                "Bridging multi-constraint spatial optimization with real-time disruption simulation", font_size=16, color=GOLD)
    
    # Left Content Column (Bullets)
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf2 = left_box.text_frame
    tf2.word_wrap = True
    
    points = [
        ("🧭 Tailored Constraint Satisfaction", "Dynamically filters and prioritizes geographic points-of-interest based on budget tiers, personal interest weights, wheelchair accessibility requirements, and dietary restrictions."),
        ("🗺️ Haversine Coordinate Optimization", "Resolves traveling route efficiency starting from hotel coords, optimizing lunch/dinner stops nearby, and strictly validating hourly opening times."),
        ("⚡ Real-Time Active Disruption Simulator", "Simulates storms, transit delays, fatigue, and sudden attraction closures on the fly, auto-replanning remaining day slots seamlessly."),
        ("🌐 Pure Client-Side Edge Execution", "Bundled using Stlite WebAssembly to run entirely in the browser, eliminating cold starts, WebSocket limits, and server costs on Vercel.")
    ]
    
    for title, desc in points:
        p_title = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
        p_title.text = "• " + title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        p_title.space_before = Pt(12)
        
        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = DARK_GRAY
        p_desc.space_before = Pt(3)
        p_desc.level = 0
        p_desc.space_after = Pt(8)

    # Right Content Column (Insert Concept Image)
    reroute_img_path = os.path.join(workspace_dir, "dynamic_rerouting.png")
    if os.path.exists(reroute_img_path):
        # Premium Card behind the image
        add_card(slide2, Inches(7.3), Inches(1.7), Inches(5.2), Inches(5.1), bg_color=LIGHT_GRAY)
        slide2.shapes.add_picture(reroute_img_path, Inches(7.5), Inches(1.9), Inches(4.8), Inches(4.7))

    # =========================================================================
    # SLIDE 3: SPATIOTEMPORAL ENGINE (Light Theme)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, WHITE)
    
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8), 
                "Heuristic Spatiotemporal Planner", font_size=28, bold=True, color=ACCENT_BLUE)
    add_textbox(slide3, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4), 
                "How VoyageFlow designs a perfect, zero-conflict day-by-day travel sequence", font_size=16, color=GOLD)
    
    # 3 Cards for Core Columns
    card_width = Inches(3.64)
    card_height = Inches(4.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)
    
    cards_data = [
        {
            "num": "01",
            "title": "Nearest-Neighbor TSP",
            "bullets": [
                "Starts each day at hotel location coordinates.",
                "Selects the next highest interest-matched sight close by.",
                "Calculates walking/driving/transit duration and fare on the fly.",
                "Minimizes cumulative transit lag dynamically."
            ]
        },
        {
            "num": "02",
            "title": "Constraint Verification",
            "bullets": [
                "Interest Weighting: Dynamically ranks places based on user vibes.",
                "Time-Window Matching: Assures arrival + duration matches opening hours.",
                "Pace Control: Limits daily stops (Relaxed 2-3, Moderate 4-5, Packed 6+)."
            ]
        },
        {
            "num": "03",
            "title": "Dietary & Health Integrity",
            "bullets": [
                "Selects budget & dietary-compliant restaurants (Vegan, Halal, Gluten-Free).",
                "Integrates lunch (12:30 PM) and dinner stops near coordinates.",
                "Excludes inaccessible landmarks for wheelchair users."
            ]
        }
    ]
    
    for i, data in enumerate(cards_data):
        left_pos = start_left + i * (card_width + gap)
        add_card(slide3, left_pos, top_pos, card_width, card_height, bg_color=LIGHT_GRAY)
        
        # Num badge
        add_textbox(slide3, left_pos + Inches(0.2), top_pos + Inches(0.2), card_width - Inches(0.4), Inches(0.6), 
                    data["num"], font_size=24, bold=True, color=GOLD)
        
        # Title
        add_textbox(slide3, left_pos + Inches(0.2), top_pos + Inches(0.8), card_width - Inches(0.4), Inches(0.6), 
                    data["title"], font_size=18, bold=True, color=NAVY)
        
        # Bullets
        bullet_box = slide3.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(1.5), card_width - Inches(0.4), Inches(3.0))
        tf3 = bullet_box.text_frame
        tf3.word_wrap = True
        for bullet in data["bullets"]:
            p = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
            p.text = "• " + bullet
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 4: THE DYNAMIC REPLANNER (Light Theme)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, WHITE)
    
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8), 
                "Real-Time Disruption & Replanning Engine", font_size=28, bold=True, color=ACCENT_BLUE)
    add_textbox(slide4, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4), 
                "Intercepting real-world events and running delta-optimizations in milliseconds", font_size=16, color=GOLD)
    
    grid_w = Inches(5.6)
    grid_h = Inches(2.2)
    grid_gap_x = Inches(0.5)
    grid_gap_y = Inches(0.4)
    
    events = [
        {
            "icon": "🌧️ Torrential Rain Storm",
            "desc": "Scans subsequent daily slots, automatically replaces all outdoor attractions with pre-seeded indoor alternatives (or cozy coffee shelters), and recalculates routes."
        },
        {
            "icon": "⏱️ Transit & Flight Delays",
            "desc": "Shifts subsequent timelines forward. If an activity is pushed past its closing hours, it automatically drops the lowest-priority spot and stretches the remaining schedule."
        },
        {
            "icon": "🚧 Landmark Closed Alert",
            "desc": "Removes the affected attraction instantly. Searches for nearby equivalent interest matches, seamlessly re-sequencing the transit segments and budget on the fly."
        },
        {
            "icon": "🥱 Traveler Fatigue Alert",
            "desc": "Detects high exhaustion levels. Drops strenuous high-effort sights, inserts coffee rest slots, and pulls forward the hotel return time for quick recovery."
        }
    ]
    
    for i, ev in enumerate(events):
        col = i % 2
        row = i // 2
        
        gx = start_left + col * (grid_w + grid_gap_x)
        gy = top_pos + row * (grid_h + grid_gap_y)
        
        add_card(slide4, gx, gy, grid_w, grid_h, bg_color=LIGHT_GRAY)
        
        add_textbox(slide4, gx + Inches(0.3), gy + Inches(0.2), grid_w - Inches(0.6), Inches(0.5), 
                    ev["icon"], font_size=18, bold=True, color=NAVY)
        
        add_textbox(slide4, gx + Inches(0.3), gy + Inches(0.7), grid_w - Inches(0.6), Inches(1.3), 
                    ev["desc"], font_size=13, color=DARK_GRAY)

    # =========================================================================
    # SLIDE 5: VISUAL EXCELLENCE (Light Theme)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, WHITE)
    
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8), 
                "Premium Glassmorphic Interface", font_size=28, bold=True, color=ACCENT_BLUE)
    add_textbox(slide5, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4), 
                "Vibrant modern aesthetics designed to wow the traveler at first glance", font_size=16, color=GOLD)
    
    # 4 Columns of Visual Features
    feat_w = Inches(2.65)
    feat_gap = Inches(0.36)
    
    features = [
        {
            "title": "Legibility & Contrast",
            "desc": "Strictly adheres to White primary and Navy Blue secondary styling. Entirely avoids light gray text/borders on white backgrounds to ensure flawless visual accessibility."
        },
        {
            "title": "3D PyDeck Maps",
            "desc": "Integrates Mapbox Light styling overlaid with sequential orange paths, deep navy location labels, and emerald restaurant marker dots."
        },
        {
            "title": "Interactive Budgets",
            "desc": "Displays real-time custom Plotly Pie charts outlining financial resource allocation (Hotel, Dining, Activities, Transit fares)."
        },
        {
            "title": "HTML Timelines",
            "desc": "Renders beautiful scrolling custom timelines with color-coded badges, visual efforts, and structured travel and rest metrics."
        }
    ]
    
    for i, feat in enumerate(features):
        fx = start_left + i * (feat_w + feat_gap)
        add_card(slide5, fx, top_pos, feat_w, Inches(4.8), bg_color=LIGHT_GRAY)
        
        # Divider line
        div = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, fx + Inches(0.2), top_pos + Inches(0.2), feat_w - Inches(0.4), Inches(0.08))
        div.fill.solid()
        div.fill.fore_color.rgb = GOLD
        div.line.fill.background()
        
        add_textbox(slide5, fx + Inches(0.2), top_pos + Inches(0.5), feat_w - Inches(0.4), Inches(0.8), 
                    feat["title"], font_size=18, bold=True, color=NAVY)
        
        add_textbox(slide5, fx + Inches(0.2), top_pos + Inches(1.4), feat_w - Inches(0.4), Inches(3.0), 
                    feat["desc"], font_size=13, color=DARK_GRAY)

    # =========================================================================
    # SLIDE 6: SERVERLESS WEBASSEMBLY (Dark Navy Theme)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, NAVY)
    
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8), 
                "Serverless WebAssembly Architecture", font_size=28, bold=True, color=WHITE)
    add_textbox(slide6, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4), 
                "Running complex Python dashboards on the Edge without backend servers", font_size=16, color=GOLD)
    
    # 2 columns for text info
    info_box_1 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf6_1 = info_box_1.text_frame
    tf6_1.word_wrap = True
    
    bullets_1 = [
        ("The Serverless Barrier", "Standard Streamlit apps require active persistent servers and active WebSockets, which are fundamentally incompatible with Vercel's stateless, ephemeral backend functions."),
        ("The Stlite Breakthrough", "VoyageFlow bypasses this by utilizing Stlite to compile Streamlit and its heavy data suites (pandas, plotly, pydeck, pydantic) into WebAssembly (Pyodide)."),
        ("Single-File HTML Bundle", "A custom build compiler (`build_stlite.py`) compresses your complete source code structure into a single Base64 string embedded in a clean index.html.")
    ]
    
    for t, d in bullets_1:
        p1 = tf6_1.add_paragraph() if tf6_1.paragraphs[0].text else tf6_1.paragraphs[0]
        p1.text = "• " + t
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = GOLD
        p1.space_before = Pt(14)
        
        p2 = tf6_1.add_paragraph()
        p2.text = d
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(220, 230, 245)
        p2.space_before = Pt(4)
        
    info_box_2 = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf6_2 = info_box_2.text_frame
    tf6_2.word_wrap = True
    
    bullets_2 = [
        ("Vercel Routing Bypass", "Vercel's automatic Python detector tries to build app.py as an ASGI server, failing because Streamlit has no 'app' export. We bypassed this elegantly with static routing."),
        ("Explicit Builds & Ignores", "Added explicit vercel.json builders targeting only static index.html and added .vercelignore filters to isolate py scripts from Vercel's compiler."),
        ("Global CDN Performance", "Because the app is fully static, it is cached on Vercel's global Edge Network, providing instant load times, zero server cost, and zero cold starts!")
    ]
    
    for t, d in bullets_2:
        p1 = tf6_2.add_paragraph() if tf6_2.paragraphs[0].text else tf6_2.paragraphs[0]
        p1.text = "• " + t
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = GOLD
        p1.space_before = Pt(14)
        
        p2 = tf6_2.add_paragraph()
        p2.text = d
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(220, 230, 245)
        p2.space_before = Pt(4)

    # Save
    prs.save(output_path)
    print(f"Presentation saved successfully at {output_path}!")

if __name__ == "__main__":
    create_presentation()
